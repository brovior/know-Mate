"""확장자 기반 실제 파일 파싱 TextExtractor."""
import logging
from pathlib import Path

from knowmate.secure.signature import UnreadableFormatError, is_zip
from knowmate.secure.text_util import format_table

logger = logging.getLogger(__name__)

_OOXML_EXTS = {".docx", ".xlsx", ".pptx"}


class PlainReader:
    def extract(self, path: str) -> str:
        """확장자에 따라 적합한 파서로 파일을 읽어 텍스트를 반환한다.

        OOXML 확장자인데 실제 zip이 아니면(DRM 래핑·손상 등) 바로
        UnreadableFormatError를 낸다. python-docx/openpyxl/python-pptx는
        내부적으로 zip을 열려다 실패하므로, 미리 걸러 의미 없는 재시도
        (예: xlsx custom.xml 복구)를 건너뛴다. COM 경유가 가능한 환경에서는
        AutoReader가 이 지점에 도달하기 전에 COM으로 라우팅한다.
        """
        ext = Path(path).suffix.lower()
        if ext in _OOXML_EXTS and not is_zip(path):
            head_hex = self._peek_hex(path)
            raise UnreadableFormatError(
                f"OOXML 확장자({ext})이나 zip 아님(DRM/암호화·손상 추정, "
                f"앞바이트 {head_hex}): {path}"
            )
        if ext == ".docx":
            return self._read_docx(path)
        if ext == ".xlsx":
            return self._read_xlsx(path)
        if ext == ".xls":
            return self._read_xls(path)
        if ext == ".pptx":
            return self._read_pptx(path)
        if ext == ".pdf":
            return self._read_pdf(path)
        if ext in {".txt", ".md", ".log"}:
            return Path(path).read_text(encoding="utf-8", errors="replace")
        raise ValueError(f"지원하지 않는 파일 형식: {ext!r} ({path})")

    @staticmethod
    def _peek_hex(path: str) -> str:
        """진단용으로 파일 앞 4바이트를 hex 문자열로 반환한다."""
        try:
            with open(path, "rb") as f:
                return f.read(4).hex().upper()
        except OSError:
            return "????"

    def _read_docx(self, path: str) -> str:
        """python-docx로 docx를 읽어 문단·표 텍스트를 문서 순서대로 반환한다."""
        import docx  # type: ignore
        from docx.document import Document as _Document
        from docx.oxml.ns import qn
        from docx.table import Table
        from docx.text.paragraph import Paragraph

        doc = docx.Document(path)
        parts: list[str] = []
        # body 자식을 순서대로 순회 → 문단과 표가 섞인 원래 순서 보존
        for child in doc.element.body.iterchildren():
            if child.tag == qn("w:p"):
                text = Paragraph(child, doc).text
                if text.strip():
                    parts.append(text)
            elif child.tag == qn("w:tbl"):
                table_text = self._format_table(
                    [[cell.text for cell in row.cells] for row in Table(child, doc).rows]
                )
                if table_text:
                    parts.append(table_text)
        return "\n".join(parts)

    @staticmethod
    def _format_table(rows: list[list[str]]) -> str:
        """표 행렬을 ' | ' 텍스트로 변환한다 (공용 format_table에 위임)."""
        return format_table(rows)

    def _read_xlsx(self, path: str) -> str:
        """openpyxl로 xlsx 파일을 읽어 셀값을 탭 구분 텍스트로 반환한다.

        손상된 사용자 정의 속성(docProps/custom.xml)으로 load_workbook이 실패하면
        해당 파트를 제거한 사본으로 재시도해 시트 데이터를 복구한다.
        """
        import openpyxl  # type: ignore

        try:
            wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
        except Exception as exc:
            # custom.xml 손상 등 메타데이터 문제 → 정리한 사본으로 재시도(복구)
            import logging
            logging.getLogger(__name__).warning(
                "xlsx 로드 실패, custom.xml 제거 후 재시도: %s (%s)", path, exc
            )
            wb = self._load_xlsx_sanitized(path)

        lines: list[str] = []
        try:
            for ws in wb.worksheets:
                sheet_lines: list[str] = []
                for row in ws.iter_rows(values_only=True):
                    row_text = "\t".join(
                        str(cell) if cell is not None else "" for cell in row
                    )
                    if row_text.strip():
                        sheet_lines.append(row_text)
                if sheet_lines:
                    # chunker가 시트 경계를 인식할 수 있도록 헤더 삽입
                    lines.append(f"=== 시트: {ws.title} ===")
                    lines.extend(sheet_lines)
        finally:
            wb.close()
        return "\n".join(lines)

    @staticmethod
    def _load_xlsx_sanitized(path: str):
        """docProps/custom.xml과 [ContentTypes].xml 내 참조를 제거한 사본으로 워크북을 로드한다.

        custom.xml 파트만 빼면 [ContentTypes].xml에 참조가 남아 openpyxl이 다시 실패하므로
        해당 Override 엔트리도 함께 제거한다.
        """
        import io
        import re
        import zipfile
        import openpyxl  # type: ignore

        CUSTOM_XML = "docProps/custom.xml"

        buf = io.BytesIO()
        with zipfile.ZipFile(path, "r") as zin:
            with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zout:
                for item in zin.infolist():
                    if item.filename == CUSTOM_XML:
                        continue
                    data = zin.read(item.filename)
                    if item.filename == "[Content_Types].xml":
                        # custom.xml Override 엔트리 제거
                        data = re.sub(
                            rb'<Override\b[^>]*/docProps/custom\.xml[^>]*/?>',
                            b"",
                            data,
                        )
                    zout.writestr(item, data)
        buf.seek(0)
        return openpyxl.load_workbook(buf, read_only=True, data_only=True)

    def _read_xls(self, path: str) -> str:
        """xlrd로 구형 xls(BIFF) 파일을 읽어 셀값을 탭 구분 텍스트로 반환한다.

        순수 파이썬 파싱이라 Office/COM이 전혀 필요 없다 — 정상 xls 대부분을
        COM 경로 밖으로 빼서 행오버·좀비 프로세스·win32timezone 문제를 원천
        차단한다(win32timezone은 COM이 날짜 셀을 변환할 때만 필요했음). DRM
        래핑이나 손상으로 xlrd가 못 여는 파일은 예외를 그대로 전파해
        AutoReader가 COM으로 폴백하게 한다(여기서 잡지 않는다).
        """
        import xlrd  # type: ignore

        wb = xlrd.open_workbook(path)
        try:
            lines: list[str] = []
            for sheet in wb.sheets():
                sheet_lines: list[str] = []
                for row_idx in range(sheet.nrows):
                    cells = [
                        self._xlrd_cell_text(sheet.cell(row_idx, col_idx), wb.datemode)
                        for col_idx in range(sheet.ncols)
                    ]
                    row_text = "\t".join(cells)
                    if row_text.strip():
                        sheet_lines.append(row_text)
                if sheet_lines:
                    # chunker가 시트 경계를 인식할 수 있도록 헤더 삽입(openpyxl 경로와 동일 포맷)
                    lines.append(f"=== 시트: {sheet.name} ===")
                    lines.extend(sheet_lines)
            return "\n".join(lines)
        finally:
            wb.release_resources()

    @staticmethod
    def _xlrd_cell_text(cell, datemode: int) -> str:
        """xlrd 셀 값을 사람이 읽을 텍스트로 변환한다(날짜·정수·불리언 서식 정리)."""
        import xlrd  # type: ignore

        if cell.ctype == xlrd.XL_CELL_EMPTY or cell.ctype == xlrd.XL_CELL_BLANK:
            return ""
        if cell.ctype == xlrd.XL_CELL_DATE:
            try:
                from xlrd.xldate import xldate_as_datetime
                dt = xldate_as_datetime(cell.value, datemode)
                if dt.time() == dt.min.time():
                    return dt.strftime("%Y-%m-%d")
                return dt.strftime("%Y-%m-%d %H:%M:%S")
            except Exception:
                return str(cell.value)
        if cell.ctype == xlrd.XL_CELL_NUMBER:
            v = cell.value
            return str(int(v)) if v == int(v) else str(v)
        if cell.ctype == xlrd.XL_CELL_BOOLEAN:
            return "TRUE" if cell.value else "FALSE"
        if cell.ctype == xlrd.XL_CELL_ERROR:
            return ""
        return str(cell.value)

    def _read_pptx(self, path: str) -> str:
        """python-pptx로 pptx를 읽어 슬라이드별 텍스트를 반환한다.

        표(has_table)와 그룹 도형(조직도) 내부 도형까지 재귀로 추출한다.
        """
        from pptx import Presentation  # type: ignore

        prs = Presentation(path)
        slides: list[str] = []
        for slide in prs.slides:
            texts: list[str] = []
            for shape in slide.shapes:
                # 도형 1개 실패가 슬라이드·파일 전체를 멈추지 않도록 건별 방어
                # (SmartArt·OLE·미디어 등 python-pptx가 인식 못 하는 도형 대응)
                try:
                    texts.extend(self._iter_shape_texts(shape))
                except Exception as exc:
                    logger.debug("pptx 도형 건너뜀 (%s): %s", type(exc).__name__, exc)
            texts = [t for t in texts if t.strip()]
            if texts:
                slides.append("\n".join(texts))
        return "\n\n".join(slides)

    def _iter_shape_texts(self, shape) -> list[str]:
        """도형 하나에서 텍스트를 추출한다. 그룹은 재귀, 표는 셀을 펼친다."""
        from pptx.enum.shapes import MSO_SHAPE_TYPE  # type: ignore

        # shape_type 접근은 인식 못 하는 도형에서 예외를 던지므로 방어한다.
        # 실패 시 그룹이 아닌 것으로 간주하고 아래 텍스트/표 추출로 진행한다.
        try:
            is_group = shape.shape_type == MSO_SHAPE_TYPE.GROUP
        except Exception:
            is_group = False

        if is_group:
            out: list[str] = []
            for child in shape.shapes:
                out.extend(self._iter_shape_texts(child))
            return out

        # has_table / has_text_frame 는 일부 도형 타입에만 존재 → getattr 로 방어
        if getattr(shape, "has_table", False):
            table_text = self._format_table(
                [[cell.text for cell in row.cells] for row in shape.table.rows]
            )
            return [table_text] if table_text else []

        if getattr(shape, "has_text_frame", False) and shape.text_frame.text.strip():
            return [shape.text_frame.text]

        return []

    def _read_pdf(self, path: str) -> str:
        """PyMuPDF(fitz)로 pdf 파일을 읽어 페이지별 텍스트를 반환한다."""
        import fitz  # type: ignore

        pages: list[str] = []
        with fitz.open(path) as doc:
            for page in doc:
                text = page.get_text()
                if text.strip():
                    pages.append(text)
        return "\n\n".join(pages)
