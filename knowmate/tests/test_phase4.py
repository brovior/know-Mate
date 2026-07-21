"""Phase 4 보안 모듈 pytest 테스트 -- fake 모드 기준, 사외 환경 전부 통과."""
from pathlib import Path
import pytest
from knowmate.rag.embedding import EmbeddingClient, VECTOR_DIM
from knowmate.secure.crypto import FakeCryptoManager, get_crypto_manager


def _fake_embed_client() -> EmbeddingClient:
    """테스트용 fake EmbeddingClient를 반환한다."""
    return EmbeddingClient(base_url="http://localhost", host_header="embed.internal", fake=True)


class TestFakeCrypto:
    def test_encrypt_returns_same(self):
        """FakeCryptoManager.encrypt()는 입력 문자열을 그대로 반환한다."""
        mgr = FakeCryptoManager()
        assert mgr.encrypt("텍스트") == "텍스트"

    def test_decrypt_returns_same(self):
        """FakeCryptoManager.decrypt()는 입력 문자열을 그대로 반환한다."""
        mgr = FakeCryptoManager()
        assert mgr.decrypt("암호문") == "암호문"

    def test_roundtrip_restores_original(self):
        """encrypt -> decrypt 왕복 후 원문이 복원된다."""
        mgr = FakeCryptoManager()
        original = "A설비 알람 폭주 처리 절차"
        assert mgr.decrypt(mgr.encrypt(original)) == original

    def test_encrypt_empty_string(self):
        """빈 문자열도 그대로 반환한다."""
        mgr = FakeCryptoManager()
        assert mgr.encrypt("") == ""
        assert mgr.decrypt("") == ""

    def test_encrypt_multiline(self):
        """여러 줄 텍스트도 그대로 반환한다."""
        mgr = FakeCryptoManager()
        text = "첫째 줄" + chr(10) + "A" + chr(10) + "셋째 줄"
        assert mgr.encrypt(text) == text


class TestCryptoInterface:
    def test_fake_mode_returns_fake_manager(self):
        """extractor=fake이면 FakeCryptoManager를 반환한다."""
        mgr = get_crypto_manager({"extractor": "fake"})
        assert isinstance(mgr, FakeCryptoManager)

    def test_default_mode_returns_fake_manager(self):
        """extractor 키 없으면 FakeCryptoManager를 반환한다."""
        mgr = get_crypto_manager({})
        assert isinstance(mgr, FakeCryptoManager)

    def test_plain_mode_requires_win32crypt(self):
        """extractor=plain이면 win32crypt 없는 환경에서 CryptoUnavailableError."""
        try:
            import win32crypt
            pytest.skip("win32crypt가 있는 환경")
        except ImportError:
            from knowmate.secure.crypto import CryptoUnavailableError
            with pytest.raises(CryptoUnavailableError):
                get_crypto_manager({"extractor": "plain"})

    def test_auto_mode_requires_win32crypt(self):
        """extractor=auto이면 win32crypt 없는 환경에서 CryptoUnavailableError."""
        try:
            import win32crypt
            pytest.skip("win32crypt가 있는 환경")
        except ImportError:
            from knowmate.secure.crypto import CryptoUnavailableError
            with pytest.raises(CryptoUnavailableError):
                get_crypto_manager({"extractor": "auto"})

class TestIndexerWithCrypto:
    def test_indexer_accepts_crypto_param(self, tmp_path: Path):
        """FakeCryptoManager를 crypto 파라미터로 Indexer 생성이 성공한다."""
        from knowmate.rag.indexer import Indexer
        embed = _fake_embed_client()
        crypto = FakeCryptoManager()
        indexer = Indexer(db_path=tmp_path / "db", embed_client=embed, crypto=crypto)
        assert indexer is not None

    def test_indexer_default_crypto_is_fake(self, tmp_path: Path):
        """crypto 파라미터 생략 시 FakeCryptoManager가 사용된다."""
        from knowmate.rag.indexer import Indexer
        embed = _fake_embed_client()
        indexer = Indexer(db_path=tmp_path / "db", embed_client=embed)
        assert isinstance(indexer._crypto, FakeCryptoManager)

    def test_encrypt_called_on_index_file(self, tmp_path: Path):
        """index_file() 시 crypto.encrypt가 청크마다 호출된다."""
        from knowmate.rag.indexer import Indexer
        from knowmate.secure.fake_reader import FakeReader
        embed = _fake_embed_client()
        call_log = []
        class SpyCrypto(FakeCryptoManager):
            def encrypt(self, plaintext: str) -> str:
                call_log.append(plaintext)
                return super().encrypt(plaintext)
        crypto = SpyCrypto()
        indexer = Indexer(db_path=tmp_path / "db", embed_client=embed, crypto=crypto)
        text = FakeReader().extract("sample.docx")
        ids = indexer.index_file(path="C:/sample/test.docx", text=text, mtime=1000.0, scope="local")
        assert len(ids) >= 1
        assert len(call_log) == len(ids)

    def test_stored_text_equals_encrypted(self, tmp_path: Path):
        """FakeCrypto 모드에서 저장된 text는 청크 평문과 동일하다."""
        from knowmate.rag.indexer import Indexer
        from knowmate.secure.fake_reader import FakeReader
        from knowmate.rag.chunker import chunk_text
        embed = _fake_embed_client()
        crypto = FakeCryptoManager()
        indexer = Indexer(db_path=tmp_path / "db", embed_client=embed, crypto=crypto)
        text = FakeReader().extract("sample.txt")
        indexer.index_file(path="C:/sample/test.txt", text=text, mtime=1000.0, scope="local")
        df = indexer.table.to_arrow().to_pandas()
        assert len(df) >= 1
        expected_chunks = chunk_text(text, "txt")
        stored_texts = sorted(df["text"].tolist())
        expected_sorted = sorted(expected_chunks)
        assert stored_texts == expected_sorted

class TestRetrieverWithCrypto:
    def _make_stack(self, tmp_path: Path, crypto=None):
        """Indexer + Retriever 스택을 반환한다."""
        from knowmate.rag.indexer import Indexer
        from knowmate.rag.retriever import Retriever
        if crypto is None:
            crypto = FakeCryptoManager()
        embed = _fake_embed_client()
        indexer = Indexer(db_path=tmp_path / "db", embed_client=embed, crypto=crypto)
        retriever = Retriever(indexer=indexer, embed_client=embed, top_k=10, score_threshold=-1.0, crypto=crypto)
        return indexer, retriever

    def test_retriever_accepts_crypto_param(self, tmp_path: Path):
        """FakeCryptoManager를 crypto 파라미터로 Retriever 생성이 성공한다."""
        _, retriever = self._make_stack(tmp_path)
        assert isinstance(retriever._crypto, FakeCryptoManager)

    def test_retriever_default_crypto_is_fake(self, tmp_path: Path):
        """crypto 파라미터 생략 시 FakeCryptoManager가 사용된다."""
        from knowmate.rag.indexer import Indexer
        from knowmate.rag.retriever import Retriever
        embed = _fake_embed_client()
        indexer = Indexer(db_path=tmp_path / "db", embed_client=embed)
        retriever = Retriever(indexer=indexer, embed_client=embed)
        assert isinstance(retriever._crypto, FakeCryptoManager)

    def test_search_returns_decrypted_text(self, tmp_path: Path):
        """search() 결과의 text가 복호화된 값이다 (FakeCrypto -> 원문 그대로)."""
        from knowmate.secure.fake_reader import FakeReader
        crypto = FakeCryptoManager()
        indexer, retriever = self._make_stack(tmp_path, crypto=crypto)
        text = FakeReader().extract("sample.docx")
        indexer.index_file(path="C:/sample/doc.docx", text=text, mtime=1000.0, scope="local")
        results = retriever.search("알람")
        assert isinstance(results, list)
        assert len(results) >= 1
        assert all(isinstance(r["text"], str) for r in results)
        assert all(len(r["text"]) > 0 for r in results)

    def test_decrypt_called_on_search_results(self, tmp_path: Path):
        """search() 시 crypto.decrypt가 결과 청크마다 호출된다."""
        from knowmate.secure.fake_reader import FakeReader
        decrypt_log = []
        class SpyCrypto(FakeCryptoManager):
            def decrypt(self, ciphertext: str) -> str:
                decrypt_log.append(ciphertext)
                return super().decrypt(ciphertext)
        crypto = SpyCrypto()
        indexer, retriever = self._make_stack(tmp_path, crypto=crypto)
        text = FakeReader().extract("sample.docx")
        indexer.index_file(path="C:/sample/doc.docx", text=text, mtime=1000.0, scope="local")
        results = retriever.search("알람")
        assert len(results) >= 1
        assert len(decrypt_log) == len(results)


class TestComReader:
    def test_com_reader_unavailable_without_win32com(self):
        """win32com 없는 환경에서 ComReader.extract() 시 ComUnavailableError."""
        try:
            import win32com.client
            pytest.skip("win32com이 있는 환경")
        except ImportError:
            from knowmate.secure.com_reader import ComReader, ComUnavailableError
            reader = ComReader()
            with pytest.raises(ComUnavailableError):
                reader.extract("C:/dummy/file.doc")

    def test_com_reader_unsupported_ext_raises(self):
        """지원하지 않는 확장자(.docx)는 ComReader에서 ValueError."""
        from knowmate.secure.com_reader import ComReader
        reader = ComReader()
        with pytest.raises((ValueError, Exception)):
            reader.extract("C:/dummy/file.docx")


class TestAutoReader:
    def test_auto_reader_docx_uses_plain(self, tmp_path: Path):
        """AutoReader.extract()로 .docx 경로를 주면 PlainReader 경로로 간다."""
        from knowmate.secure import AutoReader
        reader = AutoReader()
        fake_docx = tmp_path / "test.docx"
        with pytest.raises(Exception):
            reader.extract(str(fake_docx))

    def test_auto_reader_doc_uses_com(self):
        """AutoReader.extract()로 .doc 경로를 주면 COM 경로로 간다."""
        from knowmate.secure import AutoReader
        reader = AutoReader()
        try:
            import win32com.client
            with pytest.raises(Exception):
                reader.extract("C:/dummy/nonexistent.doc")
        except ImportError:
            from knowmate.secure.com_reader import ComUnavailableError
            with pytest.raises(ComUnavailableError):
                reader.extract("C:/dummy/nonexistent.doc")

    def test_get_extractor_auto_returns_auto_reader(self):
        """get_extractor(auto)가 AutoReader 인스턴스를 반환한다."""
        from knowmate.secure import get_extractor, AutoReader
        reader = get_extractor("auto")
        assert isinstance(reader, AutoReader)

    def test_get_extractor_plain_returns_plain_reader(self):
        """get_extractor(plain)이 PlainReader 인스턴스를 반환한다."""
        from knowmate.secure import get_extractor
        from knowmate.secure.plain_reader import PlainReader
        reader = get_extractor("plain")
        assert isinstance(reader, PlainReader)

    def test_get_extractor_fake_returns_fake_reader(self):
        """get_extractor(fake)가 FakeReader 인스턴스를 반환한다."""
        from knowmate.secure import get_extractor
        from knowmate.secure.fake_reader import FakeReader
        reader = get_extractor("fake")
        assert isinstance(reader, FakeReader)

    def test_get_extractor_unknown_raises(self):
        """알 수 없는 모드 문자열에서 ValueError."""
        from knowmate.secure import get_extractor
        with pytest.raises(ValueError):
            get_extractor("unknown_mode")


class TestOfficeGuard:
    """사용자 Office 점유 감지 → COM 파싱 연기 (OfficeBusyError)."""

    def test_process_for_ext_mapping(self):
        """확장자별 Office 실행 파일명 매핑이 올바르다."""
        from knowmate.secure.office_guard import process_for_ext
        assert process_for_ext(".doc") == "WINWORD.EXE"
        assert process_for_ext(".DOCX") == "WINWORD.EXE"
        assert process_for_ext(".xls") == "EXCEL.EXE"
        assert process_for_ext(".pptx") == "POWERPNT.EXE"
        assert process_for_ext(".pdf") is None
        assert process_for_ext(".txt") is None

    @pytest.fixture(autouse=True)
    def _clear_owned(self):
        """각 테스트 전후로 스레드 소유 PID를 비운다(테스트 간 격리)."""
        import knowmate.secure.office_guard as og
        og.clear_owned_pids()
        yield
        og.clear_owned_pids()

    def test_busy_false_when_not_running(self, monkeypatch):
        """대상 프로세스가 안 떠 있으면 False."""
        import knowmate.secure.office_guard as og
        monkeypatch.setattr(og, "_cached_processes", lambda: [("EXPLORER.EXE", 100)])
        assert og.is_office_busy_for_ext(".doc") is False

    def test_busy_true_when_running(self, monkeypatch):
        """대상 프로세스가 (우리 소유가 아닌 채) 떠 있으면 True."""
        import knowmate.secure.office_guard as og
        monkeypatch.setattr(og, "_cached_processes", lambda: [("WINWORD.EXE", 200)])
        assert og.is_office_busy_for_ext(".doc") is True
        assert og.is_office_busy_for_ext(".docx") is True

    def test_busy_false_when_running_process_is_ours(self, monkeypatch):
        """우리가 띄운(소유) 프로세스뿐이면 점유로 보지 않는다(자기 감지 방지)."""
        import knowmate.secure.office_guard as og
        monkeypatch.setattr(og, "_cached_processes", lambda: [("WINWORD.EXE", 200)])
        og.register_owned_pids({200})
        assert og.is_office_busy_for_ext(".doc") is False

    def test_busy_true_when_external_and_owned_coexist(self, monkeypatch):
        """우리 소유 인스턴스가 있어도 사용자(외부) 인스턴스가 별도로 있으면 True."""
        import knowmate.secure.office_guard as og
        monkeypatch.setattr(
            og, "_cached_processes",
            lambda: [("WINWORD.EXE", 200), ("WINWORD.EXE", 201)],
        )
        og.register_owned_pids({200})  # 200은 우리 것, 201은 사용자 것
        assert og.is_office_busy_for_ext(".doc") is True

    def test_busy_false_for_non_office_ext_even_if_running(self, monkeypatch):
        """대상 외 확장자는 프로세스가 떠 있어도 차단하지 않는다."""
        import knowmate.secure.office_guard as og
        monkeypatch.setattr(og, "_cached_processes", lambda: [("WINWORD.EXE", 200)])
        assert og.is_office_busy_for_ext(".pdf") is False

    def test_busy_false_when_enumeration_unavailable(self, monkeypatch):
        """프로세스 열거 불가(None, 비Windows 등)면 차단하지 않는다(기존 동작 유지)."""
        import knowmate.secure.office_guard as og
        monkeypatch.setattr(og, "_cached_processes", lambda: None)
        assert og.is_office_busy_for_ext(".doc") is False

    def test_clear_owned_pids_returns_and_empties(self):
        """clear_owned_pids는 기존 소유 집합을 반환하고 비운다."""
        import knowmate.secure.office_guard as og
        og.register_owned_pids({10, 11})
        assert og.clear_owned_pids() == {10, 11}
        assert og.clear_owned_pids() == set()

    def test_terminate_owned_only_kills_live_office_pids(self, monkeypatch):
        """소유 PID 중 '지금도 Office 실행 파일인' 것만 종료한다(PID 재활용 안전)."""
        import knowmate.secure.office_guard as og
        monkeypatch.setattr(og.sys, "platform", "win32")
        # 200=우리 Word(살아있음), 201=이미 죽음(목록에 없음), 202=다른 프로세스로 재활용됨
        monkeypatch.setattr(
            og, "_enumerate_processes",
            lambda: [("WINWORD.EXE", 200), ("CHROME.EXE", 202)],
        )
        killed = []
        monkeypatch.setattr(og, "_terminate_pid", lambda pid: killed.append(pid))
        og.terminate_owned_office_processes({200, 201, 202})
        assert killed == [200]  # 201(죽음)·202(재활용)는 건드리지 않음

    def test_autoreader_raises_office_busy_before_com(self, monkeypatch):
        """Office 점유 시 AutoReader가 COM 진입 전에 OfficeBusyError를 낸다."""
        import knowmate.secure.office_guard as og
        from knowmate.secure import AutoReader
        from knowmate.secure.office_guard import OfficeBusyError
        monkeypatch.setattr(og, "is_office_busy_for_ext", lambda ext: True)
        with pytest.raises(OfficeBusyError):
            AutoReader().extract("C:/dummy/report.doc")


class TestPlainReaderUnreadableFormat:
    """OOXML 확장자인데 zip이 아닌 파일(DRM 래핑·손상 등)에 대한 안전망."""

    def test_non_zip_xlsx_raises_unreadable_immediately(self, tmp_path: Path):
        """zip 아닌 xlsx는 openpyxl 재시도 없이 즉시 UnreadableFormatError."""
        from knowmate.secure.plain_reader import PlainReader
        from knowmate.secure.signature import UnreadableFormatError

        p = tmp_path / "drm.xlsx"
        p.write_bytes(b"<## " + b"\x00" * 60)
        with pytest.raises(UnreadableFormatError, match="3C232320"):
            PlainReader().extract(str(p))

    def test_non_zip_pptx_raises_unreadable(self, tmp_path: Path):
        """zip 아닌 pptx도 동일하게 UnreadableFormatError."""
        from knowmate.secure.plain_reader import PlainReader
        from knowmate.secure.signature import UnreadableFormatError

        p = tmp_path / "drm.pptx"
        p.write_bytes(b"<## " + b"\x00" * 60)
        with pytest.raises(UnreadableFormatError):
            PlainReader().extract(str(p))

    def test_ole2_xlsx_also_raises_unreadable_via_plain(self, tmp_path: Path):
        """PlainReader 단독 사용 시 OLE2 오라벨 xlsx도 UnreadableFormatError(COM 우회 없음)."""
        from knowmate.secure.plain_reader import PlainReader
        from knowmate.secure.signature import UnreadableFormatError

        p = tmp_path / "mislabeled.xlsx"
        p.write_bytes(b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1" + b"\x00" * 16)
        with pytest.raises(UnreadableFormatError):
            PlainReader().extract(str(p))


class TestPlainReaderTables:
    """표·도형 추출 검증 (docx 표 / pptx 표 + 그룹 재귀)."""

    def test_docx_extracts_table_in_order(self, tmp_path: Path):
        """docx 문단과 표가 문서 순서대로, 표는 ' | '로 추출된다."""
        import docx
        from knowmate.secure.plain_reader import PlainReader

        d = docx.Document()
        d.add_paragraph("머리말")
        t = d.add_table(rows=2, cols=2)
        t.cell(0, 0).text = "이름"; t.cell(0, 1).text = "부서"
        t.cell(1, 0).text = "홍길동"; t.cell(1, 1).text = "생산1팀"
        d.add_paragraph("꼬리말")
        p = tmp_path / "a.docx"
        d.save(str(p))

        text = PlainReader().extract(str(p))
        assert "머리말" in text
        assert "이름 | 부서" in text
        assert "홍길동 | 생산1팀" in text
        # 순서 보존: 머리말 → 표 → 꼬리말
        assert text.index("머리말") < text.index("홍길동") < text.index("꼬리말")

    def test_pptx_extracts_table_and_textbox(self, tmp_path: Path):
        """pptx 텍스트박스와 표 셀이 모두 추출된다."""
        from pptx import Presentation
        from pptx.util import Inches
        from knowmate.secure.plain_reader import PlainReader

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        tb = slide.shapes.add_textbox(Inches(0), Inches(0), Inches(3), Inches(1))
        tb.text_frame.text = "조직도"
        table = slide.shapes.add_table(2, 2, Inches(0), Inches(2), Inches(4), Inches(1)).table
        table.cell(0, 0).text = "직급"; table.cell(0, 1).text = "이름"
        table.cell(1, 0).text = "팀장"; table.cell(1, 1).text = "김철수"
        p = tmp_path / "b.pptx"
        prs.save(str(p))

        text = PlainReader().extract(str(p))
        assert "조직도" in text
        assert "직급 | 이름" in text
        assert "팀장 | 김철수" in text


# ── 시그니처 판별 ────────────────────────────────────────────────

class TestSignature:
    def test_is_ole2_true(self, tmp_path: Path):
        from knowmate.secure.signature import is_ole2
        p = tmp_path / "x.xlsx"
        p.write_bytes(b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1rest")
        assert is_ole2(str(p)) is True

    def test_is_ole2_false_for_zip(self, tmp_path: Path):
        from knowmate.secure.signature import is_ole2, is_zip
        p = tmp_path / "x.xlsx"
        p.write_bytes(b"PK\x03\x04rest")
        assert is_ole2(str(p)) is False
        assert is_zip(str(p)) is True

    def test_is_ole2_missing_file(self):
        from knowmate.secure.signature import is_ole2
        assert is_ole2("C:/nope/missing.xlsx") is False

    def test_is_zip_false_for_drm_wrapped(self, tmp_path: Path):
        """DRM 래핑 등 임의 헤더(zip도 OLE2도 아님)는 is_zip False."""
        from knowmate.secure.signature import is_ole2, is_zip
        p = tmp_path / "drm.xlsx"
        p.write_bytes(b"<## " + b"\x00" * 60)  # 실측된 DRM 래퍼 헤더 패턴
        assert is_zip(str(p)) is False
        assert is_ole2(str(p)) is False


# ── AutoReader OLE2/DRM 폴백 라우팅 ──────────────────────────────

class TestAutoReaderOle2Fallback:
    def test_ole2_labeled_xlsx_routes_to_com(self, tmp_path, monkeypatch):
        """확장자는 .xlsx인데 실제 OLE2면 ComReader로 폴백한다."""
        import knowmate.secure.com_reader as com_mod
        from knowmate.secure import AutoReader

        captured = {}

        class _FakeCom:
            def extract(self, path):
                captured["path"] = path
                return "COM_RESULT"

        monkeypatch.setattr(com_mod, "ComReader", _FakeCom)

        p = tmp_path / "mislabeled.xlsx"
        p.write_bytes(b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1" + b"\x00" * 16)  # OLE2 magic
        assert AutoReader().extract(str(p)) == "COM_RESULT"
        assert captured["path"] == str(p)

    def test_drm_wrapped_xlsx_routes_to_com(self, tmp_path, monkeypatch):
        """확장자는 .xlsx인데 zip도 OLE2도 아닌 DRM 래핑이면 ComReader로 폴백한다.

        Office는 사내 DRM 화이트리스트 프로세스라 COM으로 열면 투명 복호화된
        내용을 읽을 수 있다 — 탐색기에서는 열리는데 우리 파서만 실패하던 DRM
        문서를 이 경로로 구제한다.
        """
        import knowmate.secure.com_reader as com_mod
        from knowmate.secure import AutoReader

        captured = {}

        class _FakeCom:
            def extract(self, path):
                captured["path"] = path
                return "COM_RESULT"

        monkeypatch.setattr(com_mod, "ComReader", _FakeCom)

        p = tmp_path / "drm.xlsx"
        p.write_bytes(b"<## " + b"\x00" * 60)  # zip도 OLE2도 아닌 임의 헤더
        assert AutoReader().extract(str(p)) == "COM_RESULT"
        assert captured["path"] == str(p)

    def test_drm_wrapped_pptx_routes_to_com(self, tmp_path, monkeypatch):
        """pptx도 동일하게 DRM 래핑이면 ComReader로 폴백한다."""
        import knowmate.secure.com_reader as com_mod
        from knowmate.secure import AutoReader

        class _FakeCom:
            def extract(self, path):
                return "COM_RESULT"

        monkeypatch.setattr(com_mod, "ComReader", _FakeCom)

        p = tmp_path / "drm.pptx"
        p.write_bytes(b"<## " + b"\x00" * 60)
        assert AutoReader().extract(str(p)) == "COM_RESULT"

    def test_real_xlsx_uses_plain(self, tmp_path, monkeypatch):
        """정상 OOXML(.xlsx)은 PlainReader 경로로 간다(COM 호출 안 함)."""
        import openpyxl
        import knowmate.secure.com_reader as com_mod
        from knowmate.secure import AutoReader

        class _BoomCom:
            def extract(self, path):
                raise AssertionError("정상 xlsx인데 COM으로 갔다")

        monkeypatch.setattr(com_mod, "ComReader", _BoomCom)

        wb = openpyxl.Workbook()
        wb.active["A1"] = "정상셀"
        p = tmp_path / "ok.xlsx"
        wb.save(str(p))
        assert "정상셀" in AutoReader().extract(str(p))


# ── xlsx 손상(custom.xml) 복구 ──────────────────────────────────

class TestXlsxRecovery:
    def test_read_xlsx_recovers_when_load_fails_once(self, tmp_path, monkeypatch):
        """첫 load_workbook이 실패해도 sanitized 사본으로 복구해 셀 데이터를 읽는다."""
        import openpyxl
        from knowmate.secure.plain_reader import PlainReader

        wb = openpyxl.Workbook()
        ws = wb.active
        ws["A1"] = "헤더"; ws["B1"] = "값"
        p = tmp_path / "data.xlsx"
        wb.save(str(p))

        real_load = openpyxl.load_workbook
        state = {"n": 0}

        def flaky_load(*args, **kwargs):
            state["n"] += 1
            if state["n"] == 1:
                # custom.xml 손상 에러를 흉내냄
                raise TypeError("StringProperty.name should be str but value is NoneType")
            return real_load(*args, **kwargs)

        monkeypatch.setattr(openpyxl, "load_workbook", flaky_load)

        text = PlainReader()._read_xlsx(str(p))
        assert "헤더" in text and "값" in text
        assert state["n"] >= 2  # 직접 로드 실패 → sanitized 재시도

    def test_load_xlsx_sanitized_strips_custom_xml(self, tmp_path):
        """_load_xlsx_sanitized가 docProps/custom.xml을 제거하고 워크북을 로드한다."""
        import zipfile
        import openpyxl
        from knowmate.secure.plain_reader import PlainReader

        wb = openpyxl.Workbook()
        wb.active["A1"] = "셀값"
        p = tmp_path / "withcustom.xlsx"
        wb.save(str(p))

        # 손상된 custom.xml 주입
        with zipfile.ZipFile(str(p), "a") as z:
            z.writestr("docProps/custom.xml", "<bad/>")

        loaded = PlainReader._load_xlsx_sanitized(str(p))
        try:
            names = loaded.sheetnames
            assert names  # 시트가 정상 로드됨
        finally:
            loaded.close()


# ── .ppt COM 도형 추출 (win32com 없이 mock 도형으로 검증) ──────────

class _FakeTextRange:
    def __init__(self, text):
        self.Text = text


class _FakeTextFrame:
    def __init__(self, text):
        self.TextRange = _FakeTextRange(text)


class _FakeShape:
    """PowerPoint COM Shape의 텍스트 관련 속성만 흉내낸다."""
    def __init__(self, *, type=1, has_table=False, has_text_frame=False,
                 text="", group_items=None, table=None):
        self.Type = type
        self.HasTable = has_table
        self.HasTextFrame = has_text_frame
        self._text = text
        self.GroupItems = group_items or []
        self.Table = table

    @property
    def TextFrame(self):
        return _FakeTextFrame(self._text)


class _FakeCount:
    def __init__(self, n):
        self.Count = n


class _FakeTable:
    """grid(행×열 문자열)로 PowerPoint Table COM을 흉내낸다 (1-indexed)."""
    def __init__(self, grid):
        self._grid = grid
        self.Rows = _FakeCount(len(grid))
        self.Columns = _FakeCount(len(grid[0]) if grid else 0)

    def Cell(self, r, c):
        cell_shape = _FakeShape(has_text_frame=True, text=self._grid[r - 1][c - 1])
        return type("_C", (), {"Shape": cell_shape})()


class TestPptShapeTexts:
    """com_reader._ppt_shape_texts 의 그룹 재귀·표·텍스트프레임 분기 검증."""

    def test_text_frame_shape(self):
        from knowmate.secure.com_reader import _ppt_shape_texts
        shp = _FakeShape(has_text_frame=True, text="제목 텍스트")
        assert _ppt_shape_texts(shp) == ["제목 텍스트"]

    def test_table_shape(self):
        from knowmate.secure.com_reader import _ppt_shape_texts
        tbl = _FakeTable([["직급", "이름"], ["팀장", "김철수"]])
        shp = _FakeShape(type=19, has_table=True, table=tbl)
        assert _ppt_shape_texts(shp) == ["직급 | 이름\n팀장 | 김철수"]

    def test_group_recurses_children(self):
        from knowmate.secure.com_reader import _ppt_shape_texts
        child_box = _FakeShape(has_text_frame=True, text="대표이사")
        child_tbl = _FakeShape(
            type=19, has_table=True, table=_FakeTable([["부서", "인원"], ["생산", "10"]])
        )
        group = _FakeShape(type=6, group_items=[child_box, child_tbl])
        assert _ppt_shape_texts(group) == ["대표이사", "부서 | 인원\n생산 | 10"]

    def test_empty_shape_returns_empty(self):
        from knowmate.secure.com_reader import _ppt_shape_texts
        assert _ppt_shape_texts(_FakeShape()) == []
